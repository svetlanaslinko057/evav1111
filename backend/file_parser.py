"""
File → text extractor for the /estimate flow.

Accepts popular office/doc/image formats from web + mobile clients and
returns a single UTF-8 text blob the estimator can treat like a pasted brief.

Supported:
  • Plain text / markdown (.txt, .md)
  • PDF                    — pypdf
  • DOCX (Open XML)        — python-docx
  • XLSX (Open XML)        — openpyxl
  • PPTX (Open XML)        — python-pptx
  • Images                 — pytesseract (eng + rus) on .png/.jpg/.jpeg/.webp
  • HEIC                   — Pillow-HEIF if available, otherwise 415

Legacy binary MS formats (.doc/.xls/.ppt) are refused with a friendly
message — we won't ship libreoffice just for a brief.

Limits:
  • 10 MB max upload (enough for a brief PDF / slide deck)
  • 8000 chars max extracted text (more than enough for an estimate)
"""
from __future__ import annotations

import io
import logging
import re
import time
from typing import Optional, Tuple

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel

logger = logging.getLogger("file_parser")

MAX_BYTES = 10 * 1024 * 1024           # 10 MB
MAX_TEXT_CHARS = 8000                  # keep the estimate prompt sane


class ParsedFile(BaseModel):
    name: str
    size: int
    mime: str
    text: str
    truncated: bool = False
    source: str  # "text", "pdf", "docx", "xlsx", "pptx", "image-ocr"


# ── extractors ───────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # Tables
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
    return "\n".join(parts)


def _extract_image_ocr(data: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
    except Exception as e:
        raise HTTPException(
            status_code=415,
            detail="Image OCR is not available on this server. Try a text file or PDF.",
        ) from e
    try:
        img = Image.open(io.BytesIO(data))
        # Tesseract handles most RGB/gray modes; convert everything else.
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        # Dual-lang OCR — English + Russian cover most briefs.
        text = pytesseract.image_to_string(img, lang="eng+rus")
        return text
    except pytesseract.TesseractNotFoundError as e:
        raise HTTPException(
            status_code=415,
            detail="Image OCR engine not installed.",
        ) from e
    except Exception as e:
        logger.warning(f"OCR failed: {e}")
        raise HTTPException(
            status_code=422,
            detail=f"Could not read text from this image: {e}",
        ) from e


# ── MIME / extension routing ─────────────────────────────────────────────────

def _dispatch(filename: str, mime: str, data: bytes) -> Tuple[str, str]:
    """Return (extracted_text, source_label)."""
    name = (filename or "").lower()
    mime = (mime or "").lower()

    # Text first — cheapest.
    if mime.startswith("text/") or name.endswith((".txt", ".md", ".markdown")):
        try:
            return data.decode("utf-8", errors="replace"), "text"
        except Exception:
            return data.decode("latin-1", errors="replace"), "text"

    if name.endswith(".pdf") or mime == "application/pdf":
        return _extract_pdf(data), "pdf"

    if name.endswith(".docx") or mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return _extract_docx(data), "docx"

    if name.endswith(".xlsx") or mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ):
        return _extract_xlsx(data), "xlsx"

    if name.endswith(".pptx") or mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return _extract_pptx(data), "pptx"

    if mime.startswith("image/") or name.endswith(
        (".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif", ".bmp", ".gif")
    ):
        return _extract_image_ocr(data), "image-ocr"

    # Legacy binary MS formats — we don't support them.
    if name.endswith((".doc", ".xls", ".ppt")):
        raise HTTPException(
            status_code=415,
            detail=(
                f"Legacy {name.rsplit('.', 1)[-1].upper()} format is not "
                "supported yet. Save as .docx / .xlsx / .pptx or PDF and try again."
            ),
        )

    raise HTTPException(
        status_code=415,
        detail=f"Unsupported file type '{mime or name}'. "
               "Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, PNG/JPG/WEBP images.",
    )


# ── router ───────────────────────────────────────────────────────────────────

router = APIRouter()


# ── Voice transcription (whisper-1) ──────────────────────────────────────────
# Used by /describe → user records voice describing their product, we
# transcribe it via the active LLM provider (admin-configured in
# /admin/integrations) and the client drops the text into the goal field
# so they can edit before submitting for estimate.

import os as _os
import tempfile as _tempfile

# Whisper accepts mp3/mp4/mpeg/mpga/m4a/wav/webm (≤25 MB).
_STT_ALLOWED_EXT = {"mp3", "mp4", "mpeg", "mpga", "m4a", "wav", "webm", "ogg"}
_STT_MAX_BYTES = 25 * 1024 * 1024


@router.post("/estimate/transcribe-voice")
async def transcribe_voice_brief(file: UploadFile = File(...)):
    """Transcribe a short voice recording so the visitor can use voice
    instead of typing on /describe. Auth-free, same envelope as /estimate.

    Returns: { text: str, provider: "openai"|"emergent", chars: int }
    """
    # Lazy import — keeps backend boot fast and avoids loading the STT
    # client when nobody uses the feature.
    from stt_service import transcribe_path, STTUnavailable

    if not file:
        raise HTTPException(status_code=400, detail="No file uploaded")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty audio file")
    if len(data) > _STT_MAX_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"Audio too large ({len(data)//(1024*1024)} MB). Max 25 MB.",
        )

    # Derive extension from filename (Whisper validates by suffix). Fall back
    # to .m4a (Expo default) so the recording still gets through.
    name = (file.filename or "recording.m4a").strip() or "recording.m4a"
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else "m4a"
    if ext not in _STT_ALLOWED_EXT:
        ext = "m4a"

    # litellm/whisper wants a real on-disk file with a recognised extension,
    # so spool the upload to a temp file and pass the path.
    tmp_path: Optional[str] = None
    try:
        with _tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(data)
            tmp_path = tmp.name

        try:
            text = await transcribe_path(tmp_path)
        except STTUnavailable as e:
            raise HTTPException(status_code=503, detail=str(e))
        except Exception as e:
            logger.exception("TRANSCRIBE-VOICE failed")
            raise HTTPException(status_code=502, detail=f"Transcription failed: {e}")

        if not text:
            raise HTTPException(
                status_code=422,
                detail="Couldn't make out any speech in this recording. Try again in a quieter spot.",
            )

        logger.info(
            f"TRANSCRIBE-VOICE: name={name!r} bytes={len(data)} ext={ext} "
            f"chars={len(text)}"
        )
        return {"text": text, "chars": len(text)}
    finally:
        if tmp_path:
            try:
                _os.unlink(tmp_path)
            except Exception:
                pass


_MAX_TEXT_CHARS_OUT = 8000  # what we cap the returned brief at


class AnalyzeUrlIn(BaseModel):
    url: str
    # Optional surface tag for telemetry. Frontend sends "visitor" from
    # /describe and "admin" from the inbox panel. Anything else (or absent)
    # is recorded as "unknown" — never blocks the call.
    surface: Optional[str] = None
    # Optional device hint ("mobile" | "desktop"). Tracked only for slice;
    # we don't dispatch on it server-side. Free-form string up to 16 chars.
    device: Optional[str] = None


class AnalyzeUrlTelemetryIn(BaseModel):
    """Body for click-side events (copy / insert-into-reply / started).
       Fire-and-forget — caller never blocks waiting for the response."""
    event: str  # copy_click | insert_into_reply_click | analyze_url_started
    url: Optional[str] = None
    surface: Optional[str] = None
    device: Optional[str] = None


@router.post("/estimate/analyze-url")
async def analyze_competitor_url(body: AnalyzeUrlIn):
    """Visitor pasted a URL of a competitor / inspiration site on /describe.

    We fetch the page, distill it, and ask the active LLM to produce a
    structured brief (markdown). Caller drops `text` into the goal textarea
    so the visitor can edit before /estimate.

    Auth-free on purpose — same as /estimate. Abuse-bounded by the 8s fetch
    timeout, single-redirect cap, and the LLM provider's rate limits.

    Errors return a structured detail object — `{kind, message, hint, detail}` —
    so the UI can show a human narrative instead of a raw exception string.

    Telemetry: every call records `analyze_url_call` with `duration_ms` and
    `success`. On success: `cache_hit` or `cache_miss` (with duration). On
    failure: `analyze_url_error` (with `error_kind` and duration). The
    `analyze_url_started` event is fired by the client BEFORE this call
    so we can measure user intent independent of network outcome.
    """
    from competitor_analyzer import (
        analyze_url,
        AnalyzerUnavailable,
        FetchError,
        classify_url_error,
        classify_fetch_error,
        log_event,
    )
    from server import db as _server_db

    started_at = time.perf_counter()

    def _ms() -> int:
        return int((time.perf_counter() - started_at) * 1000)

    # Always log the call attempt — surface is best-effort.
    await log_event(
        _server_db,
        "analyze_url_call",
        url=body.url or "",
        surface=body.surface,
        device=body.device,
    )

    try:
        result = await analyze_url(body.url, db=_server_db)
    except ValueError as e:
        err = classify_url_error(str(e))
        await log_event(
            _server_db, "analyze_url_error",
            url=body.url, surface=body.surface, device=body.device,
            error_kind=err["kind"], duration_ms=_ms(), success=False,
        )
        raise HTTPException(status_code=400, detail=err)
    except FetchError as e:
        err = classify_fetch_error(str(e))
        await log_event(
            _server_db, "analyze_url_error",
            url=body.url, surface=body.surface, device=body.device,
            error_kind=err["kind"], duration_ms=_ms(), success=False,
        )
        raise HTTPException(status_code=422, detail=err)
    except AnalyzerUnavailable as e:
        err = {
            "kind": "LLM_NOT_CONFIGURED",
            "message": "Site analysis is not configured yet.",
            "hint": "An admin must enable an LLM provider in /admin/integrations.",
            "detail": str(e),
        }
        await log_event(
            _server_db, "analyze_url_error",
            url=body.url, surface=body.surface, device=body.device,
            error_kind="LLM_NOT_CONFIGURED", duration_ms=_ms(), success=False,
        )
        raise HTTPException(status_code=503, detail=err)
    except Exception as e:
        logger.exception("ANALYZE-URL failed")
        err = {
            "kind": "INTERNAL",
            "message": "Something went wrong while analyzing this link.",
            "hint": "Try again in a moment, or try a different page.",
            "detail": str(e),
        }
        await log_event(
            _server_db, "analyze_url_error",
            url=body.url, surface=body.surface, device=body.device,
            error_kind="INTERNAL", duration_ms=_ms(), success=False,
        )
        raise HTTPException(status_code=502, detail=err)

    # Cache hit / miss + duration is the cheapest signal we get. Through
    # 3-5 days these tell us if prefetch / async-mode / queue is needed.
    cached = bool(result.get("cached"))
    await log_event(
        _server_db,
        "cache_hit" if cached else "cache_miss",
        url=result.get("url") or body.url,
        surface=body.surface,
        device=body.device,
        duration_ms=_ms(),
        success=True,
    )

    text = (result.get("summary") or "").strip()
    if len(text) > _MAX_TEXT_CHARS_OUT:
        text = text[:_MAX_TEXT_CHARS_OUT].rstrip() + " …[truncated]"

    logger.info(
        f"ANALYZE-URL: url={result['url']} title={(result.get('title') or '')[:80]!r} "
        f"provider={result.get('provider')} model={result.get('model')} chars={len(text)} cached={cached} dur_ms={_ms()}"
    )
    return {
        "url": result["url"],
        "title": result.get("title") or "",
        "text": text,
        "chars": len(text),
        "provider": result.get("provider"),
        "model": result.get("model"),
        "cached": cached,
    }


@router.post("/estimate/analyze-url/telemetry")
async def analyze_url_telemetry(body: AnalyzeUrlTelemetryIn):
    """Record a click-side event. Fire-and-forget contract — the caller
    never blocks waiting for the response. Returns `{ok: true}` on accepted,
    400 only on truly unknown event names.

    Accepted events:
      • analyze_url_started        — visitor clicked ANALYZE (before network)
      • copy_click                 — admin copied the analysis text
      • insert_into_reply_click    — admin inserted the analysis into reply
    """
    from competitor_analyzer import log_event
    from server import db as _server_db

    allowed = {"analyze_url_started", "copy_click", "insert_into_reply_click"}
    if body.event not in allowed:
        raise HTTPException(
            status_code=400,
            detail={
                "kind": "INVALID_EVENT",
                "message": "Unknown telemetry event.",
                "hint": "Expected one of: analyze_url_started, copy_click, insert_into_reply_click.",
                "detail": body.event,
            },
        )
    await log_event(
        _server_db,
        body.event,
        url=body.url or "",
        surface=body.surface,
        device=body.device,
    )
    return {"ok": True}


@router.post("/estimate/parse-file", response_model=ParsedFile)
async def parse_uploaded_brief(file: UploadFile = File(...)):
    """Read a file from multipart upload and return its extracted text.

    Auth-free on purpose — same as `/estimate`, this powers the pre-login
    "we already calculated your product" experience. Abuse is bounded by the
    10 MB cap.
    """
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file")
    if len(data) > MAX_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large ({len(data) // (1024*1024)} MB). Max 10 MB.",
        )

    text, source = _dispatch(
        filename=file.filename or "",
        mime=(file.content_type or "").split(";")[0].strip(),
        data=data,
    )

    # Normalize whitespace — collapse runs of blank lines, strip NULs, rstrip.
    text = text.replace("\x00", "").strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+\n", "\n", text)

    truncated = False
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS].rstrip() + " …[truncated]"
        truncated = True

    if not text or not text.strip():
        raise HTTPException(
            status_code=422,
            detail=(
                "Couldn't extract readable text from this file. "
                "If it's a scanned image, try a clearer version."
            ),
        )

    logger.info(
        f"PARSE-FILE: {file.filename!r} mime={file.content_type} "
        f"bytes={len(data)} source={source} chars={len(text)} "
        f"truncated={truncated}"
    )
    return ParsedFile(
        name=file.filename or "file",
        size=len(data),
        mime=(file.content_type or "application/octet-stream"),
        text=text,
        truncated=truncated,
        source=source,
    )
